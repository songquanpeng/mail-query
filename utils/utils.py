import glob
import docx2txt
from utils import database
import pandas as pd
from email.header import decode_header, make_header
from email.parser import BytesParser
from io import StringIO
from pptx import Presentation
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser


def parse(path: str) -> dict:
    with open(path, 'rb') as eml_file:
        parser = BytesParser()
        message = parser.parse(eml_file)
        subject = str(make_header(decode_header(message.get("subject"))))
        sender = str(make_header(decode_header(message.get("from"))))
        receiver = str(make_header(decode_header(message.get("to"))))
        date = message.get("date")
        subject = "" if subject is None else subject
        sender = "" if sender is None else sender
        receiver = "" if receiver is None else receiver
        date = "" if date is None else date
    mail = {
        "subject": subject,
        "from": sender,
        "to": receiver,
        "date": date,
        "content": "",
        "attachments": [],
        "path": path
    }
    content = ""
    attachments = []
    last_is_plain_text = False
    for part in message.walk():
        charset = part.get_content_charset()
        if not part.is_multipart():
            content_type = part.get_content_type()
            file_name = part.get_filename()
            if file_name:
                file_name = str(make_header(decode_header(file_name)))
                file_data = part.get_payload(decode=True)
                file_data = process_attachment(file_name, file_data)
                attachments.append({"name": file_name, "content": file_data})
            else:
                if not last_is_plain_text:
                    if content_type in ['text/plain']:
                        last_is_plain_text = True
                    content = part.get_payload(decode=True)
                    if charset:
                        content = content.decode(charset)
    content = "" if content is None else content
    mail["content"] = content
    mail["attachments"] = attachments
    return mail


def process_attachment(name: str, data: bytes) -> str:
    result = ""
    if name.endswith(".txt"):
        try:
            result = data.decode("utf-8")
        except UnicodeDecodeError:
            print("unable to decode the given text by 'utf-8'")
    else:
        temp_file_path = "./data/temp"
        with open(temp_file_path, mode='wb') as temp:
            temp.write(data)
        if name.endswith(".docx"):
            result = docx2txt.process(temp_file_path)
        elif name.endswith(".pdf"):
            output_string = StringIO()
            with open(temp_file_path, mode='rb') as pdf:
                parser = PDFParser(pdf)
                doc = PDFDocument(parser)
                resource_manager = PDFResourceManager()
                device = TextConverter(resource_manager, output_string, laparams=LAParams())
                interpreter = PDFPageInterpreter(resource_manager, device)
                for page in PDFPage.create_pages(doc):
                    interpreter.process_page(page)
            result = output_string.getvalue()
        elif name.endswith(".pptx"):
            ppt = Presentation(temp_file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        result += shape.text
        elif name.endswith(".xlsx"):
            data = pd.ExcelFile(temp_file_path)
            for sheet in data.sheet_names:
                temp = data.parse(sheet)
                result += str(temp.columns)
            result += str(data.sheet_names)
    return result


def main():
    files = glob.glob("./data/*.eml")
    database.init()
    for file in files:
        mail = parse(file)
        database.insert(mail)



if __name__ == '__main__':
    main()
